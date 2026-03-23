#!/usr/bin/env python3
"""Asset Strategy Presentation Generator"""

import io
import math
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.font_manager as fm
import numpy as np

# Register IPAGothic and set as default
for f in fm.findSystemFonts():
    if 'ipagothic' in f.lower() or 'IPAGothic' in f:
        fm.fontManager.addfont(f)
        break
matplotlib.rcParams['font.family'] = ['IPAGothic', 'DejaVu Sans']
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- Colors ---
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
GREEN  = RGBColor(0x22, 0xC5, 0x5E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xE2, 0xE8, 0xF0)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
RED    = RGBColor(0xEF, 0x44, 0x44)
DKBLUE = RGBColor(0x0F, 0x1F, 0x3F)

NAVY_HEX  = "#1E3A5F"
GREEN_HEX = "#22C55E"
WHITE_HEX = "#FFFFFF"

# --- Font setup for matplotlib ---
JP_FONT = None
for f in fm.findSystemFonts():
    if 'IPAGothic' in f or 'ipagothic' in f.lower():
        JP_FONT = fm.FontProperties(fname=f)
        break

def jp(size=12):
    """Return FontProperties for Japanese text."""
    if JP_FONT:
        fp = fm.FontProperties(fname=JP_FONT.get_file(), size=size)
        return fp
    return fm.FontProperties(size=size)

# Slide dimensions (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────

def add_bg(slide, color=NAVY):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "IPAGothic"
    return txBox

def fig_to_pptx(slide, fig, x, y, w, h):
    """Save matplotlib fig to bytes and insert into slide."""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight',
                facecolor=fig.get_facecolor(), dpi=150)
    buf.seek(0)
    slide.shapes.add_picture(buf, x, y, w, h)
    plt.close(fig)


# ─────────────────────────────────────────────
# SLIDE 0: Cover
# ─────────────────────────────────────────────
def make_cover():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, NAVY)

    # accent bar left
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GREEN)

    # large title
    add_text(slide, "インフレ時代の資産戦略",
             Inches(0.5), Inches(2.0), Inches(12), Inches(1.8),
             size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # subtitle
    add_text(slide, "家族の資産会議 2026",
             Inches(0.5), Inches(3.8), Inches(12), Inches(0.8),
             size=28, bold=False, color=GREEN, align=PP_ALIGN.CENTER)

    # date
    add_text(slide, "2026年3月",
             Inches(0.5), Inches(4.6), Inches(12), Inches(0.5),
             size=18, color=LGRAY, align=PP_ALIGN.CENTER)

    # bottom bar
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), GREEN)

make_cover()


# ─────────────────────────────────────────────
# SLIDE 1: 世帯資産内訳
# ─────────────────────────────────────────────
def make_slide1():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GREEN)
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), GREEN)

    add_text(slide, "01  私たちの資産、今どうなっている？",
             Inches(0.3), Inches(0.15), Inches(12), Inches(0.6),
             size=28, bold=True, color=WHITE)

    # --- Chart ---
    categories = ["現金・預金\n(28%)", "日本株\n(32%)", "投資信託\n(22%)", "401k\n(11%)", "暗号資産他\n(3%)", "その他\n(4%)"]
    values     = [1453, 1652, 1145, 600, 131, 250]
    colors_bar = ["#64B5F6", "#22C55E", "#FDD835", "#FF8A65", "#CE93D8", "#80CBC4"]

    fig, ax = plt.subplots(figsize=(9, 4.5), facecolor=NAVY_HEX)
    ax.set_facecolor(NAVY_HEX)

    bars = ax.barh(categories, values, color=colors_bar, height=0.6)

    for bar, val in zip(bars, values):
        ax.text(bar.get_width() + 15, bar.get_y() + bar.get_height()/2,
                f"{val:,}万円", va='center', color='white',
                fontproperties=jp(10), fontsize=11)

    ax.set_xlabel("金額（万円）", fontproperties=jp(11), color='white')
    ax.tick_params(colors='white')
    ax.spines['bottom'].set_color('#4A6FA5')
    ax.spines['left'].set_color('#4A6FA5')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    for label in ax.get_yticklabels():
        label.set_fontproperties(jp(11))
        label.set_color('white')
    for label in ax.get_xticklabels():
        label.set_color('white')
    ax.set_title("世帯金融資産合計：約5,231万円", fontproperties=jp(14),
                 color=GREEN_HEX, pad=12)

    fig_to_pptx(slide, fig, Inches(0.3), Inches(0.9), Inches(8.5), Inches(5.8))

    # side note
    note_texts = [
        "💡 現金比率28%は",
        "   推奨水準（10-15%）の",
        "   約2倍",
        "",
        "世帯合計",
        "約5,231万円",
        "",
        "うち現金",
        "1,453万円",
    ]
    y_off = Inches(1.0)
    for nt in note_texts:
        if not nt:
            y_off += Inches(0.15)
            continue
        sz = 22 if "5,231" in nt or "1,453" in nt else 14
        bold = "5,231" in nt or "1,453" in nt
        col = GREEN if "💡" in nt or "2倍" in nt else WHITE
        add_text(slide, nt, Inches(9.1), y_off, Inches(4.0), Inches(0.45),
                 size=sz, bold=bold, color=col)
        y_off += Inches(0.42)

make_slide1()


# ─────────────────────────────────────────────
# SLIDE 2: インフレシナリオテーブル
# ─────────────────────────────────────────────
def make_slide2():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GREEN)
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), GREEN)

    add_text(slide, "02  インフレが進むと、現金はどうなるか",
             Inches(0.3), Inches(0.15), Inches(12), Inches(0.6),
             size=28, bold=True, color=WHITE)

    # Table via matplotlib
    rates = [-0.5, 0, 1.0, 2.0, 3.0]
    years = [1, 3, 5, 10, 15]
    principal = 1453  # 万円

    fig, ax = plt.subplots(figsize=(11, 4.8), facecolor=NAVY_HEX)
    ax.set_facecolor(NAVY_HEX)
    ax.axis('off')

    col_labels = ["インフレ率"] + [f"{y}年後" for y in years]
    row_labels = [f"{r:+.1f}%" for r in rates]

    data = []
    for r in rates:
        row = []
        for y in years:
            val = principal * ((1 - r/100) ** y)
            row.append(f"{val:,.0f}万円")
        data.append(row)

    cell_colors = []
    for ri, r in enumerate(rates):
        row_c = []
        for y in years:
            val = principal * ((1 - r/100) ** y)
            diff = val - principal
            if diff > 0:
                row_c.append("#1B4D3E")  # dark green
            elif diff > -100:
                row_c.append("#243B55")
            elif diff > -300:
                row_c.append("#5C2D2D")
            else:
                row_c.append("#7B1818")
        cell_colors.append(row_c)

    table_data = [[rl] + d for rl, d in zip(row_labels, data)]
    col_w = [0.12] + [0.176] * len(years)

    tbl = ax.table(
        cellText=table_data,
        colLabels=col_labels,
        cellLoc='center',
        loc='center',
        colWidths=col_w
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(13)

    # Style header
    for j in range(len(col_labels)):
        cell = tbl[0, j]
        cell.set_facecolor("#0F2A4A")
        cell.set_text_props(color='white', fontproperties=jp(13), fontweight='bold')
        cell.set_edgecolor('#4A6FA5')

    # Style data cells
    for ri in range(len(rates)):
        for ci in range(len(col_labels)):
            cell = tbl[ri+1, ci]
            if ci == 0:
                r = rates[ri]
                cell.set_facecolor("#0F2A4A")
                txt_col = GREEN_HEX if r <= 0 else ("#FFD600" if r <= 1 else "#FF5252")
                cell.set_text_props(color=txt_col, fontproperties=jp(13), fontweight='bold')
            else:
                cell.set_facecolor(cell_colors[ri][ci-1])
                cell.set_text_props(color='white', fontproperties=jp(12))
            cell.set_edgecolor('#4A6FA5')

    ax.set_title("1,453万円（現在の世帯現金）の実質価値推移",
                 fontproperties=jp(14), color=GREEN_HEX, pad=15)

    fig_to_pptx(slide, fig, Inches(0.3), Inches(0.9), Inches(12.5), Inches(5.8))

    add_text(slide, "※ インフレ率2%（日銀目標）で10年後に約265万円目減り",
             Inches(0.4), Inches(6.9), Inches(12), Inches(0.4),
             size=13, color=YELLOW)

make_slide2()


# ─────────────────────────────────────────────
# SLIDE 3: 現金実質価値タイムライン
# ─────────────────────────────────────────────
def make_slide3():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GREEN)
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), GREEN)

    add_text(slide, "03  現金の実質価値、年ごとの推移",
             Inches(0.3), Inches(0.15), Inches(12), Inches(0.6),
             size=28, bold=True, color=WHITE)

    years_x = list(range(0, 16))
    principal = 1453

    scenarios = {
        "デフレ −0.5%": ([-0.5], "#22C55E"),
        "現状維持 0%":   ([0.0],  "#64B5F6"),
        "インフレ 1%":   ([1.0],  "#FDD835"),
        "インフレ 2%\n（日銀目標）": ([2.0],  "#FF8A65"),
        "インフレ 3%":   ([3.0],  "#EF4444"),
    }

    fig, ax = plt.subplots(figsize=(10, 4.5), facecolor=NAVY_HEX)
    ax.set_facecolor(NAVY_HEX)

    for label, (rates, color) in scenarios.items():
        r = rates[0]
        vals = [principal * ((1 - r/100) ** y) for y in years_x]
        ax.plot(years_x, vals, color=color, linewidth=2.5, marker='o', markersize=4, label=label)

    # Highlight 2% at 15y
    val_15 = principal * (0.98 ** 15)
    ax.annotate(f"{val_15:,.0f}万円\n（−{principal-val_15:,.0f}万円）",
                xy=(15, val_15), xytext=(11.5, val_15 - 60),
                color='#FF8A65', fontproperties=jp(10),
                arrowprops=dict(arrowstyle='->', color='#FF8A65', lw=1.5))

    ax.axhline(principal, color='white', linestyle='--', linewidth=1, alpha=0.4)
    ax.set_xlabel("年数", fontproperties=jp(12), color='white')
    ax.set_ylabel("実質価値（万円）", fontproperties=jp(12), color='white')
    ax.tick_params(colors='white')
    ax.spines['bottom'].set_color('#4A6FA5')
    ax.spines['left'].set_color('#4A6FA5')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    for label_obj in ax.get_xticklabels():
        label_obj.set_color('white')
    for label_obj in ax.get_yticklabels():
        label_obj.set_color('white')

    leg = ax.legend(prop=jp(10), facecolor='#0F2A4A', edgecolor='#4A6FA5',
                    labelcolor='white', loc='lower left')

    ax.set_title("現在の世帯現金 1,453万円の実質価値シミュレーション",
                 fontproperties=jp(13), color=GREEN_HEX, pad=12)

    fig_to_pptx(slide, fig, Inches(0.3), Inches(0.9), Inches(12.5), Inches(5.8))

make_slide3()


# ─────────────────────────────────────────────
# SLIDE 4: 現金 vs 投資信託 比較
# ─────────────────────────────────────────────
def make_slide4():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GREEN)
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), GREEN)

    add_text(slide, "04  現金と投資信託、どちらがリスクか",
             Inches(0.3), Inches(0.15), Inches(12), Inches(0.6),
             size=28, bold=True, color=WHITE)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4.5), facecolor=NAVY_HEX)

    years_x = list(range(0, 16))
    principal = 1453

    # Left: Cash (2% inflation)
    ax1.set_facecolor("#1A2F4A")
    cash_vals = [principal * (0.98 ** y) for y in years_x]
    ax1.fill_between(years_x, cash_vals, principal, alpha=0.3, color='#EF4444')
    ax1.plot(years_x, cash_vals, color='#EF4444', linewidth=3)
    ax1.axhline(principal, color='white', linestyle='--', linewidth=1, alpha=0.5)
    ax1.set_title("現金（2%インフレ想定）", fontproperties=jp(13), color='#EF4444', pad=8)
    ax1.set_xlabel("年数", fontproperties=jp(11), color='white')
    ax1.set_ylabel("万円", fontproperties=jp(11), color='white')
    ax1.tick_params(colors='white')
    for s in ['top','right']: ax1.spines[s].set_visible(False)
    for s in ['bottom','left']: ax1.spines[s].set_color('#4A6FA5')
    for lbl in ax1.get_xticklabels(): lbl.set_color('white')
    for lbl in ax1.get_yticklabels(): lbl.set_color('white')
    ax1.annotate("15年後\n1,075万円\n（−265万円）",
                 xy=(15, cash_vals[-1]), xytext=(9, 1200),
                 color='#EF4444', fontproperties=jp(11),
                 arrowprops=dict(arrowstyle='->', color='#EF4444'))

    # Right: オルカン +7%/yr
    ax2.set_facecolor("#1A2F4A")
    invest_vals = [principal * (1.07 ** y) for y in years_x]
    ax2.fill_between(years_x, principal, invest_vals, alpha=0.3, color='#22C55E')
    ax2.plot(years_x, invest_vals, color='#22C55E', linewidth=3)
    ax2.axhline(principal, color='white', linestyle='--', linewidth=1, alpha=0.5)
    ax2.set_title("オルカン（年+7%想定）", fontproperties=jp(13), color='#22C55E', pad=8)
    ax2.set_xlabel("年数", fontproperties=jp(11), color='white')
    ax2.tick_params(colors='white')
    for s in ['top','right']: ax2.spines[s].set_visible(False)
    for s in ['bottom','left']: ax2.spines[s].set_color('#4A6FA5')
    for lbl in ax2.get_xticklabels(): lbl.set_color('white')
    for lbl in ax2.get_yticklabels(): lbl.set_color('white')
    ax2.annotate(f"15年後\n{invest_vals[-1]:,.0f}万円\n（+{invest_vals[-1]-principal:,.0f}万円）",
                 xy=(15, invest_vals[-1]), xytext=(8, invest_vals[-1]*0.85),
                 color='#22C55E', fontproperties=jp(11),
                 arrowprops=dict(arrowstyle='->', color='#22C55E'))

    fig.tight_layout(pad=2)
    fig_to_pptx(slide, fig, Inches(0.3), Inches(0.9), Inches(12.5), Inches(5.8))

make_slide4()


# ─────────────────────────────────────────────
# SLIDE 5: 現金ライン提案 3案
# ─────────────────────────────────────────────
def make_slide5():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GREEN)
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), GREEN)

    add_text(slide, "05  一緒に決めたいこと：現金ラインの設定",
             Inches(0.3), Inches(0.15), Inches(12), Inches(0.6),
             size=28, bold=True, color=WHITE)

    proposals = [
        {
            "title": "案A：現状維持",
            "cash": "1,453万円",
            "ratio": "現金比率 28%",
            "desc": "変更なし\n引き続き様子見",
            "pros": "変化なし・安心感",
            "cons": "インフレ損失リスク継続",
            "color": "#64B5F6",
            "tag": "現状",
        },
        {
            "title": "案B：800万円ライン ★推奨",
            "cash": "800万円",
            "ratio": "現金比率 15%",
            "desc": "差額653万円を\n段階的に積立投資へ",
            "pros": "インフレ対策 + 安心感のバランス",
            "cons": "653万円の移動が必要",
            "color": "#22C55E",
            "tag": "推奨",
        },
        {
            "title": "案C：600万円ライン",
            "cash": "600万円",
            "ratio": "現金比率 11%",
            "desc": "差額853万円を\n投資へ積極移動",
            "pros": "最大のインフレ対策効果",
            "cons": "流動性低下・心理的負担",
            "color": "#FDD835",
            "tag": "積極",
        },
    ]

    card_w = Inches(3.8)
    card_h = Inches(5.6)
    gap = Inches(0.35)
    start_x = Inches(0.5)
    start_y = Inches(0.9)

    for i, p in enumerate(proposals):
        cx = start_x + i * (card_w + gap)
        color_rgb = RGBColor(
            int(p["color"][1:3], 16),
            int(p["color"][3:5], 16),
            int(p["color"][5:7], 16)
        )
        # card bg
        add_rect(slide, cx, start_y, card_w, card_h, DKBLUE)
        # top bar
        add_rect(slide, cx, start_y, card_w, Inches(0.08), color_rgb)

        y = start_y + Inches(0.15)
        add_text(slide, p["title"], cx + Inches(0.15), y, card_w - Inches(0.3), Inches(0.55),
                 size=16, bold=True, color=color_rgb)
        y += Inches(0.6)

        add_text(slide, p["cash"], cx + Inches(0.15), y, card_w - Inches(0.3), Inches(0.65),
                 size=32, bold=True, color=WHITE)
        y += Inches(0.7)

        add_text(slide, p["ratio"], cx + Inches(0.15), y, card_w - Inches(0.3), Inches(0.4),
                 size=15, color=color_rgb)
        y += Inches(0.5)

        add_text(slide, p["desc"], cx + Inches(0.15), y, card_w - Inches(0.3), Inches(0.7),
                 size=13, color=LGRAY)
        y += Inches(0.8)

        add_text(slide, "◎ " + p["pros"], cx + Inches(0.15), y, card_w - Inches(0.3), Inches(0.6),
                 size=12, color=GREEN)
        y += Inches(0.65)

        add_text(slide, "△ " + p["cons"], cx + Inches(0.15), y, card_w - Inches(0.3), Inches(0.6),
                 size=12, color=YELLOW)

make_slide5()


# ─────────────────────────────────────────────
# SLIDE 6: まとめ
# ─────────────────────────────────────────────
def make_slide6():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GREEN)
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), GREEN)

    add_text(slide, "06  まとめ",
             Inches(0.3), Inches(0.15), Inches(12), Inches(0.6),
             size=28, bold=True, color=WHITE)

    # Metrics row
    metrics = [
        ("世帯金融資産", "5,231万円", "夫+妻+ジュニアNISA"),
        ("現金比率", "28%", "推奨の約2倍（目安10-15%）"),
        ("10年後の\nインフレ損失", "▼265万円", "2%インフレ時"),
        ("800万円ラインで\n解放される資金", "+653万円", "投資へ段階移動"),
    ]

    card_w = Inches(3.0)
    card_h = Inches(2.0)
    gap = Inches(0.2)
    mx = Inches(0.35)
    my = Inches(0.9)

    for i, (label, val, sub) in enumerate(metrics):
        cx = mx + i * (card_w + gap)
        add_rect(slide, cx, my, card_w, card_h, DKBLUE)
        add_rect(slide, cx, my, card_w, Inches(0.06), GREEN)
        add_text(slide, label, cx + Inches(0.1), my + Inches(0.1),
                 card_w - Inches(0.2), Inches(0.55),
                 size=13, color=LGRAY)
        add_text(slide, val, cx + Inches(0.1), my + Inches(0.6),
                 card_w - Inches(0.2), Inches(0.7),
                 size=26, bold=True,
                 color=RED if "▼" in val else GREEN)
        add_text(slide, sub, cx + Inches(0.1), my + Inches(1.35),
                 card_w - Inches(0.2), Inches(0.5),
                 size=11, color=LGRAY)

    # Insights
    insights = [
        "• 現金1,453万円は安全に見えるが、インフレにより毎年実質価値が低下している",
        "• 世帯合算では現金比率28%は過剰。推奨は10〜15%（800万円ライン）",
        "• ソフトバンク株が資産の約32%を占める単一銘柄集中リスクも注意が必要",
        "• 「800万円ライン」を設定することで、インフレ対策と安心感を両立できる",
    ]

    y_off = my + card_h + Inches(0.4)
    for ins in insights:
        add_text(slide, ins, Inches(0.5), y_off, Inches(12.3), Inches(0.5),
                 size=14, color=WHITE)
        y_off += Inches(0.6)

make_slide6()


# ─────────────────────────────────────────────
# SLIDE 7: End / 次のステップ
# ─────────────────────────────────────────────
def make_end():
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, GREEN)
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), GREEN)

    add_text(slide, "次のステップ",
             Inches(0.3), Inches(0.6), Inches(12), Inches(0.8),
             size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    steps = [
        ("① 現金ラインを決める", "目標：世帯現金800万円（現在1,453万円）"),
        ("② 差額653万円を段階積立", "毎月5〜10万円ペースでオルカン等へ移動"),
        ("③ 6ヶ月後に確認", "2026年9月末に資産状況・方針を再確認"),
    ]

    y = Inches(1.7)
    for step_title, step_desc in steps:
        add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.4), DKBLUE)
        add_rect(slide, Inches(0.8), y, Inches(0.07), Inches(1.4), GREEN)
        add_text(slide, step_title, Inches(1.1), y + Inches(0.1),
                 Inches(10.8), Inches(0.6),
                 size=22, bold=True, color=GREEN)
        add_text(slide, step_desc, Inches(1.1), y + Inches(0.65),
                 Inches(10.8), Inches(0.55),
                 size=16, color=LGRAY)
        y += Inches(1.6)

    add_text(slide, "一緒に考えましょう 🏠💰",
             Inches(0.3), Inches(6.7), Inches(12.5), Inches(0.5),
             size=20, color=GREEN, align=PP_ALIGN.CENTER)

make_end()


# ─────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────
out_path = "/home/user/my-game/inflation_asset_strategy.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
